#!/usr/bin/env python3
"""
将 AI 转型分享 HTML PPT 转换为 .pptx 格式。
简化版 - 使用正则提取标题和内容，再用 python-pptx 生成 PPT。
"""

import re
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 颜色方案（从原 HTML 提取）──
C = {
    "bg_dark": RGBColor(0x10, 0x18, 0x2B),
    "slide_bg": RGBColor(0x1C, 0x26, 0x44),
    "navy_2": RGBColor(0x23, 0x2F, 0x55),
    "navy_3": RGBColor(0x15, 0x1F, 0x38),
    "cream": RGBColor(0xF0, 0xEC, 0xE3),
    "cream_2": RGBColor(0xE2, 0xDC, 0xD0),
    "gold": RGBColor(0xC8, 0xA8, 0x70),
    "gold_2": RGBColor(0xEA, 0xD4, 0xA5),
    "muted": RGBColor(0x96, 0xA1, 0xB5),
    "card_bg": RGBColor(0x24, 0x2F, 0x50),
    "accent_bg": RGBColor(0x2F, 0x3F, 0x68),
    "hl_bg": RGBColor(0x2A, 0x34, 0x55),
}

FONT_CN = "Microsoft YaHei"
FONT_EN = "Arial"


def extract_slides(html_path):
    """用正则从 HTML 中提取所有幻灯片标题和 HTML 内容"""
    with open(html_path, "r", encoding="utf-8") as f:
        content = f.read()

    # 找到 slides 数组区域
    m = re.search(r"const\s+slides\s*=\s*\[(.*?)\];", content, re.DOTALL)
    if not m:
        print("❌ 未找到 slides 数组")
        sys.exit(1)
    slides_js = m.group(1)

    # 提取所有 title
    titles = re.findall(r'title:\s*"([^"]*)"', slides_js)
    
    # 提取所有 mastheadLeft / mastheadRight
    ml = re.findall(r'mastheadLeft:\s*"([^"]*)"', slides_js)
    mr = re.findall(r'mastheadRight:\s*"([^"]*)"', slides_js)
    
    # 提取所有 body (String.raw`...`)
    bodies = re.findall(r"body:\s*String\.raw`(.*?)`\s*}", slides_js, re.DOTALL)
    
    if len(titles) != len(bodies):
        print(f"⚠ 标题数({len(titles)})与内容数({len(bodies)})不匹配")
    
    slides = []
    for i, title in enumerate(titles):
        body = bodies[i] if i < len(bodies) else ""
        slides.append({
            "title": title,
            "masthead_left": ml[i] if i < len(ml) else "",
            "masthead_right": mr[i] if i < len(mr) else "",
            "body": body.strip(),
        })
    
    return slides


def strip_html(html):
    """移除 HTML 标签，返回纯文本"""
    text = re.sub(r"<[^>]+>", "", html)
    text = re.sub(r"\s+", " ", text).strip()
    return text


def extract_content_items(body):
    """从 HTML body 中抽取有意义的文本块"""
    if not body:
        return []
    
    items = []  # [(type, text), ...]   type: h1/h2/h3/p/li/b/highlight/statement/kicker/subtitle
    
    # 提取 kicker/eyebrow
    for m in re.finditer(r'class="[^"]*(?:kicker|eyebrow)[^"]*"[^>]*>(.*?)</p>', body, re.DOTALL):
        text = strip_html(m.group(1))
        if text:
            items.append(("kicker", text))
    
    # 提取 h2 标题（行内编辑 data-editable）
    for m in re.finditer(r"<h2[^>]*>(.*?)</h2>", body, re.DOTALL):
        text = strip_html(m.group(1))
        if text:
            items.append(("heading2", text))
    
    # 提取 h3 标题
    for m in re.finditer(r"<h3[^>]*>(.*?)</h3>", body, re.DOTALL):
        text = strip_html(m.group(1))
        if text:
            items.append(("heading3", text))
    
    # 提取普通的 p（排除 kicker/eyebrow）
    for m in re.finditer(r"<p[^>]*>(.*?)</p>", body, re.DOTALL):
        p_class = re.search(r'class="([^"]*)"', m.group(0))
        if p_class and any(k in p_class.group(1) for k in ("kicker", "eyebrow", "subtitle")):
            continue
        text = strip_html(m.group(1))
        if text and len(text) > 10:
            # 检测是否是重要金句
            items.append(("paragraph", text))
    
    # 提取 subtitle 类的 p
    for m in re.finditer(r'class="[^"]*subtitle[^"]*"[^>]*>(.*?)</p>', body, re.DOTALL):
        text = strip_html(m.group(1))
        if text:
            items.append(("subtitle", text))
    
    # 提取 statement-text
    for m in re.finditer(r'class="[^"]*statement-text[^"]*"[^>]*>(.*?)</div>', body, re.DOTALL):
        text = strip_html(m.group(1))
        if text:
            items.append(("statement", text))
    
    # 提取 bottom-line / note-strip / closing-line / asset-equation
    for cls in ("bottom-line", "note-strip", "closing-line", "asset-equation"):
        for m in re.finditer(f'class="[^"]*{cls}[^"]*"[^>]*>(.*?)</div>', body, re.DOTALL):
            text = strip_html(m.group(1))
            if text:
                items.append(("highlight", text))
    
    # 提取 li
    for m in re.finditer(r"<li>(.*?)</li>", body, re.DOTALL):
        text = strip_html(m.group(1))
        if text:
            items.append(("bullet", text))
    
    # 提取 b 标签（通常是强调）
    for m in re.finditer(r"<b[^>]*>(.*?)</b>", body, re.DOTALL):
        text = strip_html(m.group(1))
        if text:
            items.append(("bold", text))
    
    # 提取 card / panel / article 中的文本
    for m in re.finditer(r'<article[^>]*class="[^"]*(?:card|panel|case-step|voice|scenario)[^"]*"[^>]*>.*?</article>', body, re.DOTALL):
        # 提取其中的 tag/num
        tag_m = re.search(r'<(?:span|b) class="(?:tag|num|icon-num)">(.*?)</(?:span|b)>', m.group(0), re.DOTALL)
        tag_text = strip_html(tag_m.group(1)) if tag_m else ""
        
        # 提取全文
        full = strip_html(m.group(0))
        if tag_text:
            # 去掉 tag 部分
            rest = full.replace(tag_text, "", 1).strip()
            items.append(("card", f"[{tag_text}] {rest}" if rest else f"[{tag_text}]"))
        elif full:
            items.append(("card", full))
    
    # 提取 div 中其他有意义的内容
    for m in re.finditer(r'class="[^"]*metric-strip[^"]*"', body):
        items.append(("metric_block", ""))  # 标记有指标
    
    # 提取 ledegr-row
    for m in re.finditer(r'class="[^"]*ledger-row[^"]*"[^>]*>.*?</div>', body, re.DOTALL):
        cells = re.findall(r"<span[^>]*>(.*?)</span>", m.group(0), re.DOTALL)
        row_text = " | ".join(strip_html(c) for c in cells)
        if row_text:
            items.append(("ledger", row_text))
    
    # 提取 table-card / maturity-row
    for m in re.finditer(r'class="[^"]*maturity-row[^"]*"[^>]*>.*?</div>', body, re.DOTALL):
        cells = re.findall(r"<span[^>]*>(.*?)</span>", m.group(0), re.DOTALL)
        row_text = "  ".join(strip_html(c) for c in cells)
        if row_text:
            items.append(("ledger", row_text))
    
    # 提取 score-guide
    for m in re.finditer(r'class="[^"]*score-guide[^"]*"[^>]*>.*?</div>', body, re.DOTALL):
        spans = re.findall(r"<span[^>]*>(.*?)</span>", m.group(0), re.DOTALL)
        for s in spans:
            text = strip_html(s)
            if text:
                items.append(("badge", text))
    
    # 提取 criteria-row
    for m in re.finditer(r'class="[^"]*criteria-row[^"]*"[^>]*>.*?</div>', body, re.DOTALL):
        spans = re.findall(r"<span[^>]*>(.*?)</span>", m.group(0), re.DOTALL)
        for s in spans:
            text = strip_html(s)
            if text:
                items.append(("badge", text))
    
    # 提取 tool-pill
    for m in re.finditer(r'class="[^"]*tool-pill[^"]*"[^>]*>(.*?)</div>', body, re.DOTALL):
        text = strip_html(m.group(1))
        if text:
            items.append(("badge", text))
    
    return items


def add_tb(slide, left, top, width, height):
    """添加文本框"""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tb.text_frame.word_wrap = True
    return tb.text_frame


def add_run(para, text, size, color, bold=False, font_name=FONT_CN):
    """在段落中添加 run"""
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    return run


def create_pptx(slides, output_path):
    """根据提取的幻灯片数据创建 PPTX"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    for idx, slide_data in enumerate(slides):
        title = slide_data["title"]
        ml = slide_data["masthead_left"]
        mr = slide_data["masthead_right"]
        body = slide_data["body"]
        items = extract_content_items(body)
        
        print(f"  [{idx+1:2d}/28] {title} — {len(items)} 个内容块")
        
        is_cover = idx == 0
        
        # ── 创建幻灯片 ──
        layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(layout)
        
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = C["bg_dark"] if is_cover else C["slide_bg"]
        
        if is_cover:
            # ── 封面 ──
            # 装饰线
            ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.5), Inches(1.5), Inches(0.05))
            ln.fill.solid(); ln.fill.fore_color.rgb = C["gold"]; ln.line.fill.background()
            
            # 标题
            tf = add_tb(slide, 0.6, 2.0, 8.5, 1.5)
            add_run(tf.paragraphs[0], "传统企业 AI 转型之道", 44, C["cream"], True)
            add_run(tf.add_paragraph(), "少走弯路", 44, C["gold"], True)
            
            # 副标题
            tf2 = add_tb(slide, 0.6, 3.8, 7, 0.8)
            add_run(tf2.paragraphs[0], '从"工具焦虑"回到"业务场景"：先找到一个能提效、能降本、能算清账的流程。', 16, C["cream_2"])
            
            # 三个标签
            for i, (n, l) in enumerate([("01","戳痛点"),("02","讲经历"),("03","落案例")]):
                sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6+i*2.8), Inches(5.0), Inches(2.4), Inches(0.55))
                sh.fill.solid(); sh.fill.fore_color.rgb = C["navy_2"]
                sh.line.color.rgb = C["gold"]; sh.line.width = Pt(0.5)
                tf_i = sh.text_frame; tf_i.vertical_anchor = MSO_ANCHOR.MIDDLE
                tf_i.paragraphs[0].alignment = PP_ALIGN.CENTER
                add_run(tf_i.paragraphs[0], f"{n} · {l}", 14, C["gold"], True)
            
            # 演讲人线
            ln2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(6.0), Inches(0.8), Inches(0.03))
            ln2.fill.solid(); ln2.fill.fore_color.rgb = C["gold"]; ln2.line.fill.background()
            
            tf3 = add_tb(slide, 1.6, 5.85, 4, 0.5)
            add_run(tf3.paragraphs[0], "肖洒", 20, C["cream"], True)
            add_run(tf3.add_paragraph(), "原银行金融业高级 AI 产品经理 · 现传统企业 AI 转型负责人", 11, C["muted"])
            
            # 底部标签
            bdg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(7.0), Inches(2.2), Inches(0.4))
            bdg.fill.solid(); bdg.fill.fore_color.rgb = C["navy_2"]
            bdg.line.color.rgb = C["gold"]; bdg.line.width = Pt(0.5)
            bdg.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            bdg.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            add_run(bdg.text_frame.paragraphs[0], "HOOK / STORY / ROI", 9, C["muted"], font_name=FONT_EN)
        else:
            # ── 普通页 ──
            # 左上角标签
            if ml:
                tf_l = add_tb(slide, 0.6, 0.2, 4, 0.4)
                add_run(tf_l.paragraphs[0], ml, 10, C["gold"], False, FONT_EN)
            
            # 右上角
            if mr:
                tf_r = add_tb(slide, 6, 0.2, 3.5, 0.4)
                tf_r.paragraphs[0].alignment = PP_ALIGN.RIGHT
                add_run(tf_r.paragraphs[0], mr, 10, C["muted"], False, FONT_EN)
            
            # 页面标题
            tf_t = add_tb(slide, 0.6, 0.75, 8.8, 0.7)
            add_run(tf_t.paragraphs[0], title, 26, C["cream"], True)
            
            # 分隔线
            ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.55), Inches(1.0), Inches(0.04))
            ln.fill.solid(); ln.fill.fore_color.rgb = C["gold"]; ln.line.fill.background()
            
            # ── 内容区 ──
            y = 1.8
            max_y = 7.0
            col_w = 8.6
            
            # 先预分类
            cards = [it for it in items if it[0] in ("card", "ledger", "metric_block")]
            others = [it for it in items if it[0] not in ("card", "ledger", "metric_block")]
            has_many_cards = len(cards) >= 2
            
            if has_many_cards:
                # 两栏卡片排列
                rendered_cards = []
                for item in items:
                    if y > max_y:
                        break
                    t, txt = item
                    
                    if t == "kicker":
                        tf_k = add_tb(slide, 0.6, y, col_w, 0.35)
                        add_run(tf_k.paragraphs[0], txt, 11, C["gold"], True, FONT_EN)
                        y += 0.3
                    
                    elif t == "card":
                        card_w = 4.15
                        card_h = 1.3
                        idx_c = len(rendered_cards)
                        col = idx_c % 2
                        row = idx_c // 2
                        cx = 0.6 + col * (card_w + 0.3)
                        cy = y + row * (card_h + 0.12)
                        
                        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx), Inches(cy), Inches(card_w), Inches(card_h))
                        sh.fill.solid(); sh.fill.fore_color.rgb = C["card_bg"]
                        sh.line.color.rgb = C["navy_2"]; sh.line.width = Pt(0.5)
                        sh.text_frame.word_wrap = True
                        sh.text_frame.margin_left = Inches(0.15); sh.text_frame.margin_right = Inches(0.1)
                        sh.text_frame.margin_top = Inches(0.1)
                        
                        add_run(sh.text_frame.paragraphs[0], txt[:250], 10, C["cream_2"], True)
                        rendered_cards.append(item)
                    
                    elif t == "bullet":
                        tf_b = add_tb(slide, 0.8, y, col_w-0.3, 0.3)
                        add_run(tf_b.paragraphs[0], f"• {txt[:200]}", 9, C["cream_2"])
                        y += 0.28
                    
                    elif t == "bold":
                        tf_bo = add_tb(slide, 0.8, y, col_w-0.3, 0.3)
                        add_run(tf_bo.paragraphs[0], f"→ {txt[:200]}", 10, C["gold_2"], True)
                        y += 0.28
                    
                    elif t == "highlight":
                        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(col_w+0.3), Inches(0.5))
                        sh.fill.solid(); sh.fill.fore_color.rgb = C["hl_bg"]
                        sh.line.color.rgb = C["gold"]; sh.line.width = Pt(1.2)
                        sh.text_frame.word_wrap = True; sh.text_frame.margin_left = Inches(0.15)
                        sh.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                        add_run(sh.text_frame.paragraphs[0], txt[:200], 10, C["gold_2"], True)
                        y += 0.52
                    
                    elif t == "statement":
                        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(col_w+0.3), Inches(0.7))
                        sh.fill.solid(); sh.fill.fore_color.rgb = C["hl_bg"]
                        sh.line.color.rgb = C["navy_2"]; sh.line.width = Pt(0.5)
                        sh.text_frame.word_wrap = True; sh.text_frame.margin_left = Inches(0.2)
                        sh.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                        sh.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                        add_run(sh.text_frame.paragraphs[0], txt[:300], 13, C["cream"], True)
                        y += 0.75
                    
                    elif t == "ledger":
                        tf_lg = add_tb(slide, 0.6, y, col_w, 0.25)
                        add_run(tf_lg.paragraphs[0], txt[:300], 8, C["cream_2"])
                        y += 0.22
                    
                    elif t == "badge":
                        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(y), Inches(2.0), Inches(0.3))
                        sh.fill.solid(); sh.fill.fore_color.rgb = C["navy_2"]
                        sh.line.color.rgb = C["gold"]; sh.line.width = Pt(0.3)
                        sh.text_frame.word_wrap = True; sh.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                        sh.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                        add_run(sh.text_frame.paragraphs[0], txt[:40], 8, C["gold_2"], True, FONT_EN)
                        y += 0.35
                    
                    elif t in ("heading2", "heading3"):
                        sz = 16 if t == "heading2" else 13
                        tf_h = add_tb(slide, 0.6, y, col_w, 0.4)
                        add_run(tf_h.paragraphs[0], txt[:200], sz, C["cream"], True)
                        y += 0.38
                    
                    elif t == "paragraph":
                        tf_p = add_tb(slide, 0.6, y, col_w, 0.4)
                        add_run(tf_p.paragraphs[0], txt[:300], 11, C["cream_2"])
                        y += 0.38
                    
                    elif t == "subtitle":
                        tf_s = add_tb(slide, 0.6, y, col_w, 0.5)
                        add_run(tf_s.paragraphs[0], txt[:300], 14, C["cream_2"])
                        y += 0.45
                
                # 卡片后的 y 更新（动态计算）
                num_rows = (len(rendered_cards) + 1) // 2
                y = max(y, y + num_rows * 1.42)
            else:
                # 单栏
                for item in items:
                    if y > max_y:
                        break
                    t, txt = item
                    
                    if t == "kicker":
                        tf_k = add_tb(slide, 0.6, y, col_w, 0.35)
                        add_run(tf_k.paragraphs[0], txt, 11, C["gold"], True, FONT_EN)
                        y += 0.3
                    
                    elif t in ("heading2", "heading3"):
                        sz = 16 if t == "heading2" else 13
                        tf_h = add_tb(slide, 0.6, y, col_w, 0.4)
                        add_run(tf_h.paragraphs[0], txt[:200], sz, C["cream"], True)
                        y += 0.38
                    
                    elif t == "paragraph":
                        tf_p = add_tb(slide, 0.6, y, col_w, 0.4)
                        add_run(tf_p.paragraphs[0], txt[:300], 11, C["cream_2"])
                        y += 0.38
                    
                    elif t == "subtitle":
                        tf_s = add_tb(slide, 0.6, y, col_w, 0.5)
                        add_run(tf_s.paragraphs[0], txt[:300], 14, C["cream_2"])
                        y += 0.45
                    
                    elif t == "card":
                        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(col_w+0.3), Inches(0.6))
                        sh.fill.solid(); sh.fill.fore_color.rgb = C["card_bg"]
                        sh.line.color.rgb = C["navy_2"]; sh.line.width = Pt(0.5)
                        sh.text_frame.word_wrap = True; sh.text_frame.margin_left = Inches(0.15)
                        sh.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                        add_run(sh.text_frame.paragraphs[0], txt[:250], 10, C["cream_2"], True)
                        y += 0.62
                    
                    elif t == "bullet":
                        tf_b = add_tb(slide, 0.8, y, col_w-0.3, 0.3)
                        add_run(tf_b.paragraphs[0], f"• {txt[:200]}", 9, C["cream_2"])
                        y += 0.28
                    
                    elif t == "bold":
                        tf_bo = add_tb(slide, 0.8, y, col_w-0.3, 0.3)
                        add_run(tf_bo.paragraphs[0], f"→ {txt[:200]}", 10, C["gold_2"], True)
                        y += 0.28
                    
                    elif t == "highlight":
                        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(col_w+0.3), Inches(0.5))
                        sh.fill.solid(); sh.fill.fore_color.rgb = C["hl_bg"]
                        sh.line.color.rgb = C["gold"]; sh.line.width = Pt(1.2)
                        sh.text_frame.word_wrap = True; sh.text_frame.margin_left = Inches(0.15)
                        sh.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                        add_run(sh.text_frame.paragraphs[0], txt[:200], 10, C["gold_2"], True)
                        y += 0.52
                    
                    elif t == "statement":
                        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(col_w+0.3), Inches(0.7))
                        sh.fill.solid(); sh.fill.fore_color.rgb = C["hl_bg"]
                        sh.line.color.rgb = C["navy_2"]; sh.line.width = Pt(0.5)
                        sh.text_frame.word_wrap = True; sh.text_frame.margin_left = Inches(0.2)
                        sh.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                        sh.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                        add_run(sh.text_frame.paragraphs[0], txt[:300], 13, C["cream"], True)
                        y += 0.75
                    
                    elif t == "ledger":
                        tf_lg = add_tb(slide, 0.6, y, col_w, 0.25)
                        add_run(tf_lg.paragraphs[0], txt[:300], 8, C["cream_2"])
                        y += 0.22
                    
                    elif t == "badge":
                        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(y), Inches(2.0), Inches(0.3))
                        sh.fill.solid(); sh.fill.fore_color.rgb = C["navy_2"]
                        sh.line.color.rgb = C["gold"]; sh.line.width = Pt(0.3)
                        sh.text_frame.word_wrap = True; sh.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                        sh.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                        add_run(sh.text_frame.paragraphs[0], txt[:40], 8, C["gold_2"], True, FONT_EN)
                        y += 0.35
            
            # ── 页码 ──
            total = len(slides)
            tf_page = add_tb(slide, 0.6, 7.05, 3, 0.3)
            add_run(tf_page.paragraphs[0], "Enterprise AI Playbook 2.0", 8, C["muted"], font_name=FONT_EN)
            
            tf_pg = add_tb(slide, 7.5, 7.05, 2, 0.3)
            tf_pg.paragraphs[0].alignment = PP_ALIGN.RIGHT
            add_run(tf_pg.paragraphs[0], f"{idx+1:02d} / {total:02d}", 8, C["muted"], font_name=FONT_EN)
    
    prs.save(output_path)
    print(f"\n✅ 保存成功: {output_path}")


def main():
    input_html = "/Users/xiaohan/Documents/金融/ai-share-ppt/index.html"
    output_pptx = "/Users/xiaohan/Documents/金融/ai-share-ppt/AI转型分享.pptx"
    
    if len(sys.argv) > 1:
        input_html = sys.argv[1]
    if len(sys.argv) > 2:
        output_pptx = sys.argv[2]
    
    print(f"📖 读取: {input_html}")
    slides = extract_slides(input_html)
    print(f"   → 共 {len(slides)} 张幻灯片\n")
    
    print("🔄 正在生成 PPTX...")
    create_pptx(slides, output_pptx)
    
    file_size = __import__("os").path.getsize(output_pptx)
    print(f"📦 文件大小: {file_size/1024:.0f} KB")


if __name__ == "__main__":
    main()
